from pptx import Presentation
import re
import os
import csv
import argparse

def extract_slide_titles(filename):
    """
    Extracts slide titles and slide numbers from a PowerPoint file, skipping hidden slides.
    :param filename: Path to the PowerPoint file.
    :return: A list of dictionaries containing filenames, slide numbers, and titles.
    """
    prs = Presentation(filename)
    slide_data = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        # Check if slide is hidden
        if hasattr(slide, "element") and slide.element.get("show") == "0":
            continue

        title = slide.shapes.title
        if title:
            title_text = re.sub(r'\s+', ' ', title.text.strip())
            slide_data.append({
                "filename": filename,
                "slide_number": slide_index,
                "slide_title": title_text
            })

    return slide_data

def process_directory(directory):
    """
    Finds all PowerPoint files in a directory and extracts their slide titles.
    :param directory: Path to the directory.
    :return: A list of dictionaries with filenames, slide numbers, and titles.
    """
    pptx_files = [os.path.join(directory, f) for f in os.listdir(directory) if f.endswith('.pptx')]
    results = []
    for pptx_file in pptx_files:
        results.extend(extract_slide_titles(pptx_file))
    return results

def write_to_csv(output_file, data):
    """
    Writes extracted slide titles and slide numbers to a CSV file.
    :param output_file: Path to the output CSV file.
    :param data: List of dictionaries with filenames, slide numbers, and titles.
    """
    with open(output_file, mode='w', newline='', encoding='utf-8') as csv_file:
        writer = csv.writer(csv_file)
        writer.writerow(["Filename", "Slide Number", "Slide Title"])

        for item in data:
            writer.writerow([item["filename"], item["slide_number"], item["slide_title"]])

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract slide titles and numbers from PowerPoint files.")
    parser.add_argument("-f", "--file", help="Path to a single PowerPoint file.")
    parser.add_argument("-d", "--directory", help="Path to a directory containing PowerPoint files.")
    parser.add_argument("-o", "--output", default="output.csv", help="Path to the output CSV file.")
    args = parser.parse_args()

    if args.file:
        # Process a single file
        data = extract_slide_titles(args.file)
    elif args.directory:
        # Process a directory of files
        data = process_directory(args.directory)
    else:
        raise ValueError("You must provide either a file (-f) or a directory (-d).")

    # Write results to a CSV file
    write_to_csv(args.output, data)
    print(f"Output written to {args.output}")